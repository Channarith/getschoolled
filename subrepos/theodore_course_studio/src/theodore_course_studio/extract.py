"""Extract page/slide text from PDF and PPTX for training + course generation."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass
class ExtractedPage:
    index: int
    title: str
    text: str
    marked_reject_hint: bool = False  # best-effort from annotations when available


@dataclass
class ExtractedDocument:
    path: str
    pages: list[ExtractedPage]
    extractor: str
    error: str = ""


def extract_document(path: str | Path) -> ExtractedDocument:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    return ExtractedDocument(path=str(path), pages=[], extractor="none", error=f"unsupported:{suffix}")


def _extract_pdf(path: Path) -> ExtractedDocument:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as exc:  # noqa: BLE001 — optional dep
        return ExtractedDocument(
            path=str(path),
            pages=[],
            extractor="pypdf-missing",
            error=f"pypdf unavailable: {exc}",
        )
    try:
        reader = PdfReader(str(path))
        pages: list[ExtractedPage] = []
        for i, page in enumerate(reader.pages):
            try:
                text = (page.extract_text() or "").strip()
            except Exception:  # noqa: BLE001
                text = ""
            # Ink/annotation reject heuristic: free-text annots mentioning reject/no
            # or presence of many ink annotations — humans used circle+line marks.
            marked = False
            try:
                annots = page.get("/Annots") or []
                if annots and len(list(annots)) >= 2:
                    # Dense annotations often mean reviewer markup on the page.
                    marked = True
            except Exception:  # noqa: BLE001
                marked = False
            title = text.splitlines()[0][:120] if text else f"Page {i + 1}"
            pages.append(
                ExtractedPage(index=i, title=title, text=text, marked_reject_hint=marked)
            )
        return ExtractedDocument(path=str(path), pages=pages, extractor="pypdf")
    except Exception as exc:  # noqa: BLE001
        return ExtractedDocument(path=str(path), pages=[], extractor="pypdf", error=str(exc))


def _extract_pptx(path: Path) -> ExtractedDocument:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception as exc:  # noqa: BLE001
        return ExtractedDocument(
            path=str(path),
            pages=[],
            extractor="python-pptx-missing",
            error=f"python-pptx unavailable: {exc}",
        )
    try:
        prs = Presentation(str(path))
        pages: list[ExtractedPage] = []
        for i, slide in enumerate(prs.slides):
            chunks: list[str] = []
            marked = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text.strip())
                # Circle + line-through style rejects often appear as freeform /
                # line / oval drawings left by the reviewer.
                try:
                    st = shape.shape_type
                    if st in {
                        MSO_SHAPE_TYPE.AUTO_SHAPE,
                        MSO_SHAPE_TYPE.LINE,
                        MSO_SHAPE_TYPE.FREEFORM,
                    }:
                        name = (getattr(shape, "name", "") or "").lower()
                        if any(tok in name for tok in ("oval", "ellipse", "circle", "line", "scribble")):
                            marked = True
                except Exception:  # noqa: BLE001
                    pass
            text = "\n".join(c for c in chunks if c)
            title = chunks[0][:120] if chunks else f"Slide {i + 1}"
            pages.append(
                ExtractedPage(index=i, title=title, text=text, marked_reject_hint=marked)
            )
        return ExtractedDocument(path=str(path), pages=pages, extractor="python-pptx")
    except Exception as exc:  # noqa: BLE001
        return ExtractedDocument(path=str(path), pages=[], extractor="python-pptx", error=str(exc))
