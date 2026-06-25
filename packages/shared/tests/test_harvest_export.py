"""Part 1 -> Part 2 bridge: export a generated course to .pptx + course.json."""

import json

import pytest

from aoep_shared.harvest import (
    CourseTags,
    export_course_package,
    extract_text,
    generate_course,
)

_SAMPLE = (
    "Introduction\nWelcome to algebra; we cover the core objectives.\n\n"
    "Example 1\nA worked example solving for x.\n\n"
    "Summary\nIn summary, algebra solves for unknowns.\n"
)


def _course():
    doc = extract_text(_SAMPLE, default_title="Algebra 101")
    return generate_course(doc, subject="math", tags=CourseTags(core_fundamental=True))


def test_course_json_export_roundtrips(tmp_path):
    pkg = export_course_package(_course(), tmp_path, write_pptx=False)
    assert pkg.course_json_path and pkg.course_json_path.exists()
    data = json.loads(pkg.course_json_path.read_text())
    assert data["title"] == "Algebra 101"
    assert data["composition_score"] == pkg.composition_score
    assert "core-fundamental" in data["tags"]["labels"]


def test_pptx_export_is_readable_by_pptx(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    course = _course()
    pkg = export_course_package(course, tmp_path, write_pptx=True)
    assert pkg.pptx_path and pkg.pptx_path.exists()

    prs = Presentation(str(pkg.pptx_path))
    # title slide + one per generated slide
    assert len(prs.slides) == len(course.slides) + 1
    # The first content slide title matches the first generated slide.
    content_titles = [s.shapes.title.text for s in list(prs.slides)[1:] if s.shapes.title]
    assert course.slides[0].title in content_titles
    # Narration is carried in speaker notes (what ppt_trainer reads as notes).
    notes_blob = " ".join(
        s.notes_slide.notes_text_frame.text
        for s in list(prs.slides)[1:] if s.has_notes_slide
    )
    assert course.slides[0].narration.split(".")[0] in notes_blob
