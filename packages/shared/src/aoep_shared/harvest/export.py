"""Export a generated course to formats the teaching layer (Part 2) consumes.

Part 1 (the harvester) produces a ``GeneratedCourse``. Part 2 (the agentic
ppt_trainer reader) natively reads ``.pptx``/``.pdf``. This module is the bridge:

  - ``export_pptx``          -> a .pptx the ppt_trainer reads with zero changes
                               (one slide per generated slide: title + bullet
                               body, with the harvester narration carried in the
                               slide's SPEAKER NOTES so Part 2 can reuse it).
  - ``export_course_json``    -> the rich course package (slides + composition
                               matrix + PCS score + tags) for the meeting layer.
  - ``export_course_package`` -> writes both into one directory and returns a
                               manifest (the hand-off contract between parts).

python-pptx is required for every harvest export (Part 1 always writes .pptx).
"""

from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional

from .generate import GeneratedCourse

PPTX_INSTALL_HINT = (
    "python-pptx is required for harvest export. Install with: "
    "pip install -e 'packages/shared[harvest]'  "
    "(or: pip install python-pptx==0.6.23)"
)


def ensure_pptx_available() -> None:
    """Raise ImportError with install instructions when python-pptx is missing."""
    try:
        import pptx  # noqa: F401
    except ImportError as exc:
        raise ImportError(PPTX_INSTALL_HINT) from exc


def resolve_course_pptx(course_json: str | Path) -> Optional[Path]:
    """Resolve the sibling ``.pptx`` for a ``*.course.json`` harvest package."""
    path = Path(course_json)
    if not path.name.endswith(".course.json"):
        sibling = path.with_suffix(".pptx")
        return sibling if sibling.is_file() else None
    candidate = path.with_name(path.name[: -len(".course.json")] + ".pptx")
    if candidate.is_file():
        return candidate
    pptxs = sorted(path.parent.glob("*.pptx"))
    return pptxs[0] if len(pptxs) == 1 else None


def _bullets(body: str, *, max_bullets: int = 8) -> List[str]:
    """Split slide body into on-screen bullet lines (newline-aware)."""
    lines = [ln.strip() for ln in (body or "").splitlines() if ln.strip()]
    if lines:
        return lines[:max_bullets]
    parts: List[str] = []
    for chunk in body.replace("\n", " ").split(". "):
        s = chunk.strip().rstrip(".")
        if s:
            parts.append(s if len(s) <= 140 else s[:137] + "...")
    return parts[:max_bullets] or [body[:140]]


def export_pptx(course: GeneratedCourse, path: str | Path) -> Path:
    """Write the course as a .pptx (title slide + one slide per generated slide).

    The spoken narration the harvester already produced is stored in each slide's
    speaker notes, so Part 2 (ppt_trainer) can read it back as ``Section.notes``.
    """
    from pptx import Presentation  # lazy

    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # Title slide.
    title_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_layout)
    s0.shapes.title.text = course.title
    if len(s0.placeholders) > 1:
        subtitle = f"Interactive lesson · {course.subject}"
        if course.tags and course.tags.career_path:
            subtitle += f" · path: {course.tags.career_path}"
        s0.placeholders[1].text = subtitle

    # Content slides.
    content_layout = prs.slide_layouts[1]
    for slide in course.slides:
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = slide.title
        body_tf = s.placeholders[1].text_frame
        bullets = _bullets(slide.body)
        body_tf.text = bullets[0]
        for line in bullets[1:]:
            p = body_tf.add_paragraph()
            p.text = line
            p.level = 1 if line.startswith((
                "More depth:", "Did you know?", "Hands-on:", "Code walkthrough:",
                "Worked example", "From the text:", "Step ",
            )) else 0
        if slide.media_url:
            p = body_tf.add_paragraph()
            p.text = f"▶ Watch demo: {slide.media_url}"
            p.level = 0
        notes_parts = []
        if slide.narration:
            notes_parts.append(slide.narration)
        if slide.audio_path:
            notes_parts.append(f"[Audio narration: {slide.audio_path}]")
        if slide.media_url:
            kind = slide.media_kind or "media"
            notes_parts.append(f"[{kind.title()} demo: {slide.media_url}]")
        if slide.category:
            notes_parts.append(f"[Slide type: {slide.category}]")
        if len(slide.body) > 400:
            notes_parts.append(f"[Full reference text]\n{slide.body}")
        if notes_parts:
            notes = s.notes_slide.notes_text_frame
            notes.text = "\n\n".join(notes_parts)

    prs.save(str(path))
    return path


def export_course_json(course: GeneratedCourse, path: str | Path, *, theme: Optional[dict] = None) -> Path:
    """Write the rich course package JSON (slides + composition + score + tags)."""
    path = Path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = course.to_dict()
    if theme:
        payload["theme"] = theme
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return path


@dataclass
class CoursePackage:
    """Hand-off contract from Part 1 (harvest) to Parts 2 (teach) / 3 (present)."""

    course_id: str
    title: str
    subject: str
    composition_score: int
    output_dir: Path
    pptx_path: Optional[Path] = None
    course_json_path: Optional[Path] = None

    def to_dict(self) -> Dict:
        return {
            "course_id": self.course_id,
            "title": self.title,
            "subject": self.subject,
            "composition_score": self.composition_score,
            "output_dir": str(self.output_dir),
            "pptx_path": str(self.pptx_path) if self.pptx_path else None,
            "course_json_path": str(self.course_json_path) if self.course_json_path else None,
        }


def export_course_package(
    course: GeneratedCourse,
    out_dir: str | Path,
    *,
    write_pptx: bool = True,
    with_media: bool = False,
    repo_root: Optional[Path] = None,
) -> CoursePackage:
    """Write the course package (pptx + json) into ``out_dir`` and return a manifest."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    stem = (course.title or course.course_id or "course").strip().replace("/", "-")[:64] or "course"

    if with_media:
        from .media import export_course_media
        export_course_media(course, out_dir, repo_root=repo_root)

    from .themes import resolve_slide_theme
    theme = resolve_slide_theme(
        title=course.title,
        subject=course.subject,
        tags=course.tags.label_list() if course.tags else (),
        fmt=course.fmt,
    ).to_dict()

    json_path = export_course_json(course, out_dir / f"{stem}.course.json", theme=theme)
    pptx_path: Optional[Path] = None
    if write_pptx:
        ensure_pptx_available()
        pptx_path = export_pptx(course, out_dir / f"{stem}.pptx")

    return CoursePackage(
        course_id=course.course_id,
        title=course.title,
        subject=course.subject,
        composition_score=course.composition_score,
        output_dir=out_dir,
        pptx_path=pptx_path,
        course_json_path=json_path,
    )
