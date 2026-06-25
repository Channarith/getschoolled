"""Multi-source content extractors for the harvester.

Reads teaching material from many source types and normalizes each into an
``ExtractedDoc`` (title + ordered ``(heading, text)`` sections) that the
generator turns into a course. One normalized shape means the composition
matrix / scoring / tagging downstream are source-agnostic.

Supported source types:
  - "text"        raw/markdown-ish text (stdlib)
  - "html"/"url"  web pages                      (beautifulsoup4, lazy)
  - "pdf"         PDF documents                  (pypdf, lazy)
  - "pptx"        PowerPoint decks               (python-pptx, lazy)
  - "docx"        Word documents                 (python-docx, lazy)
  - "database"    rows from a SQL database        (stdlib sqlite3; any DB-API
                  connection can be passed in for Postgres/MySQL/etc.)

Heavy parsers are imported lazily so importing this module (and the worker)
never requires the optional deps; only the extractor you call needs its lib.
"""

from __future__ import annotations

import re
import sqlite3
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

SUPPORTED_SOURCE_TYPES: Tuple[str, ...] = (
    "text", "html", "url", "pdf", "pptx", "docx", "database",
)

_HEADING_HINT = re.compile(r"^(chapter|section|unit|lesson|part|module)\b", re.IGNORECASE)


@dataclass
class ExtractedDoc:
    """Normalized extraction result shared by every source type."""

    title: str
    language: str = "en"
    source_type: str = "text"
    sections: List[Tuple[str, str]] = field(default_factory=list)  # (heading, text)

    def nonempty_sections(self) -> List[Tuple[str, str]]:
        return [(h, t) for h, t in self.sections if (t or "").strip()]


def _looks_like_heading(line: str) -> bool:
    s = line.strip()
    if not s or len(s) > 80:
        return False
    if _HEADING_HINT.match(s):
        return True
    if s[-1] in ".,;:":
        return False
    words = s.split()
    return 1 <= len(words) <= 10 and (s.isupper() or s.istitle())


# --------------------------------------------------------------------------- #
# text
# --------------------------------------------------------------------------- #
def extract_text(text: str, *, default_title: str = "Untitled") -> ExtractedDoc:
    """Split plain text into sections on blank-line/heading boundaries."""
    title = default_title
    sections: List[Tuple[str, str]] = []
    heading = title
    buf: List[str] = []

    def flush() -> None:
        nonlocal buf
        if buf:
            sections.append((heading, " ".join(buf).strip()))
            buf = []

    lines = text.splitlines()
    for raw in lines:
        line = raw.strip()
        if not line:
            flush()
            continue
        if line.startswith("#"):  # markdown-ish heading
            flush()
            heading = line.lstrip("# ").strip() or heading
            continue
        if _looks_like_heading(line):
            flush()
            heading = line
            continue
        buf.append(line)
    flush()
    return ExtractedDoc(title=title, source_type="text",
                        sections=[s for s in sections if s[1]])


# --------------------------------------------------------------------------- #
# html / url
# --------------------------------------------------------------------------- #
def extract_html(html: str, *, default_title: str = "Untitled") -> ExtractedDoc:
    from bs4 import BeautifulSoup  # lazy

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    title = default_title
    if soup.title and soup.title.string:
        title = soup.title.string.strip()
    elif soup.find("h1"):
        title = soup.find("h1").get_text(strip=True)

    sections: List[Tuple[str, str]] = []
    heading = title
    body: List[str] = []
    for el in soup.find_all(["h1", "h2", "h3", "p", "li"]):
        txt = el.get_text(" ", strip=True)
        if not txt:
            continue
        if el.name in ("h1", "h2", "h3"):
            if body:
                sections.append((heading, " ".join(body).strip()))
                body = []
            heading = txt
        else:
            body.append(txt)
    if body:
        sections.append((heading, " ".join(body).strip()))
    if not sections:
        whole = soup.get_text(" ", strip=True)
        if whole:
            sections = [(title, whole)]
    return ExtractedDoc(title=title, source_type="html",
                        sections=[s for s in sections if s[1]])


# --------------------------------------------------------------------------- #
# pdf
# --------------------------------------------------------------------------- #
def extract_pdf(data: bytes, *, default_title: str = "Untitled") -> ExtractedDoc:
    from pypdf import PdfReader  # lazy

    reader = PdfReader(BytesIO(data))
    title = default_title
    meta_title = getattr(reader.metadata, "title", None) if reader.metadata else None
    if meta_title:
        title = str(meta_title).strip() or default_title

    sections: List[Tuple[str, str]] = []
    heading: Optional[str] = None
    buf: List[str] = []

    def flush() -> None:
        nonlocal buf
        if buf:
            sections.append((heading or title, " ".join(buf).strip()))
            buf = []

    for page in reader.pages:
        for raw in (page.extract_text() or "").splitlines():
            line = raw.strip()
            if not line:
                continue
            if _looks_like_heading(line):
                flush()
                heading = line
            else:
                buf.append(line)
    flush()
    return ExtractedDoc(title=title, source_type="pdf",
                        sections=[s for s in sections if s[1]])


# --------------------------------------------------------------------------- #
# pptx (PowerPoint)
# --------------------------------------------------------------------------- #
def extract_pptx(data: bytes, *, default_title: str = "Untitled") -> ExtractedDoc:
    """One section per slide: the slide title becomes the heading, the remaining
    text frames become the body."""
    from pptx import Presentation  # lazy (python-pptx)

    prs = Presentation(BytesIO(data))
    sections: List[Tuple[str, str]] = []
    title = default_title
    for idx, slide in enumerate(prs.slides):
        slide_title = ""
        body_parts: List[str] = []
        # The title placeholder, when present, is the slide heading.
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = slide.shapes.title.text.strip()
        except (AttributeError, ValueError):
            slide_title = ""
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            txt = shape.text_frame.text.strip()
            if not txt or txt == slide_title:
                continue
            body_parts.append(txt)
        heading = slide_title or f"Slide {idx + 1}"
        if idx == 0 and slide_title:
            title = slide_title
        sections.append((heading, " ".join(body_parts).strip()))
    return ExtractedDoc(title=title, source_type="pptx",
                        sections=[s for s in sections if s[1]])


# --------------------------------------------------------------------------- #
# docx (Word)
# --------------------------------------------------------------------------- #
def extract_docx(data: bytes, *, default_title: str = "Untitled") -> ExtractedDoc:
    """Use Word heading styles (Heading 1/2/3, Title) as section boundaries."""
    from docx import Document  # lazy (python-docx)

    doc = Document(BytesIO(data))
    title = default_title
    sections: List[Tuple[str, str]] = []
    heading: Optional[str] = None
    buf: List[str] = []

    def flush() -> None:
        nonlocal buf
        if buf:
            sections.append((heading or title, " ".join(buf).strip()))
            buf = []

    for para in doc.paragraphs:
        txt = (para.text or "").strip()
        if not txt:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.lower().startswith("title") and title == default_title:
            title = txt
            heading = txt
            continue
        if style.lower().startswith("heading"):
            flush()
            heading = txt
            continue
        buf.append(txt)
    flush()
    return ExtractedDoc(title=title, source_type="docx",
                        sections=[s for s in sections if s[1]])


# --------------------------------------------------------------------------- #
# database
# --------------------------------------------------------------------------- #
def extract_database(
    *,
    connection=None,
    db_path: Optional[str] = None,
    query: str,
    heading_column: Optional[str] = None,
    text_columns: Optional[Sequence[str]] = None,
    title: str = "Database course",
) -> ExtractedDoc:
    """Read rows from a SQL database into sections (one section per row).

    Pass either a live DB-API ``connection`` (sqlite3/psycopg2/mysql-connector)
    or a sqlite ``db_path``. ``heading_column`` supplies each section heading
    (falls back to the row index); ``text_columns`` selects the body columns
    (defaults to every non-heading column).
    """
    own = False
    if connection is None:
        if not db_path:
            raise ValueError("provide either connection or db_path")
        connection = sqlite3.connect(db_path)
        own = True
    try:
        cur = connection.cursor()
        cur.execute(query)
        col_names = [d[0] for d in (cur.description or [])]
        rows = cur.fetchall()
    finally:
        if own:
            connection.close()

    sections: List[Tuple[str, str]] = []
    for r_idx, row in enumerate(rows):
        record = dict(zip(col_names, row))
        heading = ""
        if heading_column and heading_column in record:
            heading = str(record[heading_column])
        heading = heading or f"Row {r_idx + 1}"
        cols = text_columns or [c for c in col_names if c != heading_column]
        body = " ".join(str(record[c]) for c in cols if record.get(c) is not None).strip()
        sections.append((heading, body))
    return ExtractedDoc(title=title, source_type="database",
                        sections=[s for s in sections if s[1]])


# --------------------------------------------------------------------------- #
# dispatch
# --------------------------------------------------------------------------- #
def extract(source_type: str, data, *, default_title: str = "Untitled", **kwargs) -> ExtractedDoc:
    """Dispatch to the right extractor by source type.

    ``data`` is text/bytes for file/web types and ignored for "database"
    (pass connection/db_path/query as kwargs).
    """
    st = (source_type or "text").lower()
    if st == "text":
        return extract_text(_as_text(data), default_title=default_title)
    if st in ("html", "url"):
        return extract_html(_as_text(data), default_title=default_title)
    if st == "pdf":
        return extract_pdf(_as_bytes(data), default_title=default_title)
    if st == "pptx":
        return extract_pptx(_as_bytes(data), default_title=default_title)
    if st == "docx":
        return extract_docx(_as_bytes(data), default_title=default_title)
    if st == "database":
        return extract_database(title=default_title, **kwargs)
    raise ValueError(f"unsupported source_type {source_type!r}; one of {SUPPORTED_SOURCE_TYPES}")


def extract_file(path: str, *, source_type: Optional[str] = None,
                 default_title: Optional[str] = None) -> ExtractedDoc:
    """Extract from a local file, inferring the source type from the extension."""
    p = Path(path)
    st = source_type or _infer_source_type(p)
    title = default_title or p.stem
    if st in ("text", "html", "url"):
        return extract(st, p.read_text(encoding="utf-8", errors="replace"), default_title=title)
    return extract(st, p.read_bytes(), default_title=title)


def _infer_source_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    return {
        "pdf": "pdf", "pptx": "pptx", "ppt": "pptx",
        "docx": "docx", "doc": "docx",
        "html": "html", "htm": "html",
        "txt": "text", "md": "text", "rst": "text",
        "db": "database", "sqlite": "database", "sqlite3": "database",
    }.get(ext, "text")


def _as_text(data) -> str:
    if isinstance(data, bytes):
        return data.decode("utf-8", errors="replace")
    return str(data)


def _as_bytes(data) -> bytes:
    if isinstance(data, str):
        return data.encode("utf-8")
    return bytes(data)
